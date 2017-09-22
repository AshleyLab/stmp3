#a script that exports excel sheets created by genetic counselors to powerpoint slides
#author Noah Friedman mail4noahf@gmail.com
from pptx import *
from pptx.util import *
from pptx.enum.text import *
from pptx.enum.shapes import *
from pptx.dml.color import *
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import sys

import xls_parsing_functions

#################################################################################
#constants for layout
#Size of the pptx slide
#ALERT I should programmatically access the real value so this stays stable
PPTX_WIDTH = Inches(10.0)
PPTX_HEIGHT = Inches(7.5)

#how far stuff on the left is offset
LEFT_OFFSET = Inches(0.3)

#Formatting for the tables
BETWEEN_TABLE_OFFSET = Inches(0.3)
#parameter only used for initialization
TABLE_WIDTH = Inches(2.75)
#these lists are a list of (displayTableName, excelColumnName) pairs
#Alert the RVIS value is not validated
IN_SILICO_TABLE_ROW_NAMES = [('SIFT:', 'SIFT Function'),
('PolyPhen:', 'PolyPhen-2 Function'),
('MutationTaster:', 'MutationTaster'),
('RVIS:', 'RVIS'), ('CADD:', 'CADD Score'),
('PhylopP100:', 'phyloP100way'),
('UCSC:', 'UCSC')]
#Alert the exac pop max value is incorrect
ALLELE_TABLE_ROW_NAMES = [('ExAC (overall):','ExAC (%)'),
('ExAC (popmax):', '?'),
('gnomAD (overall):', 'gnomad'),
('gnomAD (popmax):', 'gnomad'),
('1000Genomes:', '1000 Genomes')]
ROWS_IN_SILICO_TABLE = len(IN_SILICO_TABLE_ROW_NAMES) + 1
ROWS_ALLELE_TABLE = len(ALLELE_TABLE_ROW_NAMES) + 1
TABLE_COLUMN_WIDTH = Inches(1.5)
TABLE_ROW_HEIGHT = Inches(0.25)
TABLE_TEXT_SIZE = 8

#left offset for text boxes on the right hand of the screen
TEXT_BOX_LEFT_OFFSET = TABLE_WIDTH + LEFT_OFFSET + Inches(.5)
TEXT_BOX_WIDTH = Inches(6.25)
#DEFAULT_TEXTBOX_HEIGHT = Inches(1.75)
#ALERT THIS IS ARTIFICIALLY LOW FOR THE sake of me being able to see what im doing
DEFAULT_TEXTBOX_HEIGHT = Inches(1)
#constants for the title boxes
TITLE_TOP_OFFSET = Inches(0.1)
TITLE_HEIGHT = Inches(0.75)
TITLE_WIDTH = Inches(5)
#offsets for the top of the tables/text boxes underneath the title
ELEMENT_TOP_OFFSET = TITLE_HEIGHT + TITLE_TOP_OFFSET + Inches(.5)
OFFSET_BETWEEN_STUDIES_BOX = Inches(.25) 

#formatting for the udn logo
LOGO_OFFSET_FROM_SIDE = Inches(0.25)
#constants for the png we use
UDN_PNG_WIDTH = 508
UDN_PNG_HEIGHT = 388
UDN_PNG_HEIGHT_TO_WIDTH_RATIO = 1.0*UDN_PNG_HEIGHT/UDN_PNG_WIDTH
UDN_LOGO_WIDTH = Inches(1.0)
UDN_LOGO_HEIGHT = UDN_LOGO_WIDTH*UDN_PNG_HEIGHT_TO_WIDTH_RATIO

#constants for displaying the udn id at the top
UDN_ID_TOP_OFFSET = Inches(0)
UDN_ID_LEFT_OFFSET = Inches(8.0)

##################################################################################

#Helper Methods for making slides

#sets the text size and performs magic with text size, font etc
def set_text_size_and_font(shape, size):
	tf = shape.text_frame
	p = tf.paragraphs[0]
	p.font.size = Pt(size)

#makes text in a paragraph italics, bold and underlined if needed, adjusts size if specified
def apply_italics_bold_underlining_and_size(p, text, italic, bold, underlined, size = Pt(10)):
	run = p.add_run()
	run.text = text
	font = run.font
	if italic: font.italic = True
	if bold: font.bold = True
	if underlined: font.underlined = True

#sets the shapes color
def set_shape_color(shape, r, g, b):
	fill = shape.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(r, g, b)

#sets the shapes text while making the label italics and the rest normal
def set_shape_text(shape, labelText, sourceText):
	text_frame = shape.text_frame
	text_frame.clear()
	p = text_frame.paragraphs[0]
	p.text = labelText + sourceText
	p.font.italic = True 
	p.font.size = Pt(10)
	p.alignment = PP_ALIGN.LEFT

###########################################################################################
#Add the UDN logo to the corner
def display_udn_logo(slide):
	#convert them to inches
	logoLeft = PPTX_WIDTH - UDN_LOGO_WIDTH - LOGO_OFFSET_FROM_SIDE
	logoTop = PPTX_HEIGHT - UDN_LOGO_HEIGHT - LOGO_OFFSET_FROM_SIDE
	slide.shapes.add_picture('UDN_logo_temp_screenshot.png', logoLeft, logoTop, UDN_LOGO_WIDTH, UDN_LOGO_HEIGHT)

#Add the UDN id to the corner
def display_udn_id(slide, udnId):
	txBox = slide.shapes.add_textbox(UDN_ID_LEFT_OFFSET, UDN_ID_TOP_OFFSET, Inches(1.0), Inches(1.0))
	txBox.text = udnId

#creates the gene name display at the top of the slide
def display_gene_name(topOfSlideTextbox, dfRow):
	geneName = xls_parsing_functions.get_xls_value(dfRow, 'Gene Symbol')
	topOfSlideTextbox.text = "GENE: " + geneName
	#we get the first paragraph (which is set by deafult), and make it bold
	topOfSlideTextbox.paragraphs[0].font.bold = True

#displays the nm, c and p below the gene name
def display_tscriptId_tscriptVariant_proteinVariant_and_exon(topOfSlideTextbox, dfRow):
	tscriptId = xls_parsing_functions.get_xls_value(dfRow,'Transcript ID', '****')
	tscriptVariant = xls_parsing_functions.get_xls_value(dfRow,'Transcript Variant', '****')
	proteinVariant = xls_parsing_functions.get_xls_value(dfRow,'Protein Variant', '****')
	#exon = xls_parsing_functions.get_xls_value(dfRow,'Protein Variant', '****')
	#alert please change 
	exon = 'exon'
	text = ' | '.join([tscriptId, tscriptVariant, proteinVariant, exon])
	#add the text to the slide 
	p = topOfSlideTextbox.add_paragraph()
	p.text = text

def display_chr_pos_ref_alt(topOfSlideTextbox, dfRow):
	chromosome = xls_parsing_functions.get_xls_value(dfRow,'Chromosome', '****')
	pos = xls_parsing_functions.get_xls_value(dfRow,'Position', '****')
	ref = xls_parsing_functions.get_xls_value(dfRow,'Reference Allele', '****')
	alt = xls_parsing_functions.get_xls_value(dfRow,'Sample Allele', '****')
	chrPos = ':'.join([chromosome, pos])
	refAlt = ">".join([ref,alt])
	text = ''.join([chrPos, refAlt])
	#add the text to the slide
	p = topOfSlideTextbox.add_paragraph()
	p.text = text

#displays the textbox that we show on the top of the slide (this include chr pos, transcript variant etc)
def display_top_of_slide_textbox(slide, dfRow):
	txBox = slide.shapes.add_textbox(LEFT_OFFSET, TITLE_TOP_OFFSET, TITLE_WIDTH, TITLE_HEIGHT)
	tf = txBox.text_frame
	display_gene_name(tf, dfRow)
	display_tscriptId_tscriptVariant_proteinVariant_and_exon(tf, dfRow)
	display_chr_pos_ref_alt(tf, dfRow)

##########################################################################################
#this function returns the amount of space it used to display its text
def display_gene_function_textbox(slide, dfRow):
	#ALERT--add logic that resizes the text boxes based on how much text there is
	if(True): geneFunctionHeight = DEFAULT_TEXTBOX_HEIGHT
	shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TEXT_BOX_LEFT_OFFSET, ELEMENT_TOP_OFFSET, TEXT_BOX_WIDTH, geneFunctionHeight)
	set_shape_color(shape, 80, 172, 196)
	info = xls_parsing_functions.get_xls_value(dfRow,'Gene Function')
	set_shape_text(shape, "Gene function: ", info)
	return ELEMENT_TOP_OFFSET + geneFunctionHeight

#this function returns the amount of space it used to display its text plus the amount the gene_function textbox used
def display_functional_studies_textbox(slide, dfRow, currentOffsetFromTop):
	#if text in the gene function box is bleeding over, make sure to resize my new text box so it is lower down
	functionalStudiesOffset = currentOffsetFromTop + OFFSET_BETWEEN_STUDIES_BOX	
	#ALERT--add logic that resizes the text boxes based on how much text there is
	if(True): functionalStudiesHeight = DEFAULT_TEXTBOX_HEIGHT
	shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TEXT_BOX_LEFT_OFFSET, functionalStudiesOffset, TEXT_BOX_WIDTH, functionalStudiesHeight)
	set_shape_color(shape, 156, 186, 95)
	info = xls_parsing_functions.get_xls_value(dfRow, 'Animal Model')
	set_shape_text(shape, "Functional studies: ", info)
	return currentOffsetFromTop + DEFAULT_TEXTBOX_HEIGHT + OFFSET_BETWEEN_STUDIES_BOX

def display_known_disease_association_textbox(slide, dfRow, currentOffsetFromTop):
	knownDiseaseAssocOffset = currentOffsetFromTop + OFFSET_BETWEEN_STUDIES_BOX
	#ALERT--add logic that resizes the text boxes based on how much text there is
	if(True): knownDiseaseAssocHeight = DEFAULT_TEXTBOX_HEIGHT
	shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TEXT_BOX_LEFT_OFFSET, knownDiseaseAssocOffset, TEXT_BOX_WIDTH, knownDiseaseAssocHeight)
	set_shape_color(shape, 190, 7, 18)
	#Alert we need to change this
	info = xls_parsing_functions.get_xls_value(dfRow, 'Clinical Features')
	set_shape_text(shape, "Known Disease Association: ", info)

######################################################################################################
#given a table and the coordinates of a cell changes its fill color to fill color
def adjust_cell_color(table, i, j, fillColor):
	c = table.cell(i, j).fill.solid()
	table.cell(i, j).fill.fore_color.rgb = fillColor

#sets the headers of the tables we use
def set_table_header_rows(slide, table, rightHeader, leftHeader):
	adjust_cell_color(table, 0, 0, RGBColor(0, 0, 0))
	adjust_cell_color(table, 0, 1, RGBColor(0, 0, 0))
	set_text_size_and_font(table.cell(0, 0), TABLE_TEXT_SIZE)
	set_text_size_and_font(table.cell(0, 1), TABLE_TEXT_SIZE)
	table.cell(0, 0).text = rightHeader
	table.cell(0, 1).text = leftHeader

#set the width and height dimensions of a table
def set_table_dimensions(table, nRows):
	#note all our table have two columns and an indeterminate number of rows
	table.columns[0].width = TABLE_COLUMN_WIDTH
	table.columns[1].width = TABLE_COLUMN_WIDTH
	#set all the rows
	for i in range(nRows): table.rows[i].height = TABLE_ROW_HEIGHT

#sets each row of the tables we use
#now that rowNames is a list of tuples: (displayRowName, excelKey), in the order we want to display
def set_table_body_rows(slide, table, dfRow, rowNames):
	idx = 1
	for displayRowName, excelKey in rowNames:
		rowValue = xls_parsing_functions.get_xls_value(dfRow, excelKey)
		set_text_size_and_font(table.cell(idx, 0), TABLE_TEXT_SIZE)
		set_text_size_and_font(table.cell(idx, 1), TABLE_TEXT_SIZE)
		table.cell(idx, 0).text = displayRowName
		table.cell(idx, 1).text = rowValue
		#Make the cells white
		adjust_cell_color(table, idx, 0, RGBColor(255, 255, 255))
		adjust_cell_color(table, idx, 1, RGBColor(255, 255, 255))
		idx += 1


def display_in_silico_prediction_table(slide, dfRow):
	#add the table to the slide
	modelTable = slide.shapes.add_table(ROWS_IN_SILICO_TABLE, 2, LEFT_OFFSET, ELEMENT_TOP_OFFSET, TABLE_WIDTH, TABLE_ROW_HEIGHT*ROWS_IN_SILICO_TABLE)
	#note that to get the actual table object we must deal with the table aspect of the table object
	modelTable = modelTable.table
	set_table_header_rows(slide, modelTable, 'Program', 'Prediction')
	set_table_body_rows(slide, modelTable, dfRow, IN_SILICO_TABLE_ROW_NAMES)
	set_table_dimensions(modelTable, len(IN_SILICO_TABLE_ROW_NAMES))
	return ELEMENT_TOP_OFFSET + TABLE_ROW_HEIGHT*ROWS_IN_SILICO_TABLE
	#maybe we want this idiomatic code??
	#modelTable.ApplyStyle('3C2FFA5D-87B4-456A-9821-1D502468CF0F')
	#modelTable.first_row = True
	
def display_allele_frequency_table(slide, dfRow, currentOffsetFromTop):
	currentOffsetFromTop += BETWEEN_TABLE_OFFSET
	alleleFreqTable = slide.shapes.add_table(ROWS_ALLELE_TABLE, 2, LEFT_OFFSET, currentOffsetFromTop, TABLE_WIDTH, TABLE_ROW_HEIGHT*ROWS_ALLELE_TABLE)
	alleleFreqTable = alleleFreqTable.table
	set_table_dimensions(alleleFreqTable, len(ALLELE_TABLE_ROW_NAMES))
	set_table_header_rows(slide, alleleFreqTable, 'Source & Ethnicity', 'MAF %')
	set_table_body_rows(slide, alleleFreqTable, dfRow, ALLELE_TABLE_ROW_NAMES)
	set_table_dimensions(alleleFreqTable, len(ALLELE_TABLE_ROW_NAMES))
	return currentOffsetFromTop + TABLE_ROW_HEIGHT*ROWS_ALLELE_TABLE

#used to create the IGV and variant deicsion tables
def display_1d_table(slide, currentOffsetFromTop, headerText, fillText):
	currentOffsetFromTop += BETWEEN_TABLE_OFFSET
	oneDResultTable = slide.shapes.add_table(1, 2, LEFT_OFFSET, currentOffsetFromTop, TABLE_WIDTH, TABLE_ROW_HEIGHT*ROWS_ALLELE_TABLE)
	oneDResultTable = oneDResultTable.table
	#set the text first because we are going to be messing with the colors
	apply_italics_bold_underlining_and_size(oneDResultTable.cell(0, 0).text_frame.paragraphs[0], headerText, True, False, False)
	#ALERt it for some reason forces us to have white even thoigh we dont want white
	apply_italics_bold_underlining_and_size(oneDResultTable.cell(0, 1).text_frame.paragraphs[0], fillText, False, False, False)
	#For IGV one cell is black and the other is white
	adjust_cell_color(oneDResultTable, 0, 0, RGBColor(0, 0, 0))
	adjust_cell_color(oneDResultTable, 0, 1, RGBColor(255, 255, 255))
	set_text_size_and_font(oneDResultTable.cell(0, 0), TABLE_TEXT_SIZE)
	set_text_size_and_font(oneDResultTable.cell(0, 1), TABLE_TEXT_SIZE)
	set_table_dimensions(oneDResultTable, 1)

	return currentOffsetFromTop + TABLE_ROW_HEIGHT


#given a blank template slide and a variant json record this code fully populates a template slide
def populate_slide(slide, dfRow, udnId):

	####################################
	display_udn_logo(slide)
	display_udn_id(slide, udnId)
	#CREATE THE TOP OF THE SLIDE: Gene, name, chrom, pos etc
	#we use current offset from the top as an intermediate variable that helps us put stuff in the right place
	display_top_of_slide_textbox(slide, dfRow)
	currentOffsetFromTop = display_gene_function_textbox(slide, dfRow)
	currentOffsetFromTop = display_functional_studies_textbox(slide, dfRow, currentOffsetFromTop)
	display_known_disease_association_textbox(slide, dfRow, currentOffsetFromTop)
	currentOffsetFromTop = display_in_silico_prediction_table(slide, dfRow)
	currentOffsetFromTop = display_allele_frequency_table(slide, dfRow, currentOffsetFromTop)
	#Do the 1-D tables--note that in the interest of minimal copy pasting of code some code from the display allele freq tables is re
	igvStatus = xls_parsing_functions.get_xls_value(dfRow,'IGV')
	currentOffsetFromTop = display_1d_table(slide, currentOffsetFromTop, 'IGV', igvStatus)
	display_1d_table(slide, currentOffsetFromTop, 'Variant Decision', 'tbd')

######################################################################################################
print 'the file should be run by typing the command in this exact manner: '
print 'python powerpoint_export.py excelSpreadsheetFilename UDNID'
print 'the excel spreadsheet must either be in the same directory as this piece of code or you must explicitly specify its filepath'
print 'make sure your excel filename doesnt have spaces! Spaces confuse linux'
print 'beginning export'

#command line parsing 
xlsFilename = sys.argv[1]
udnId = sys.argv[2]

#we build up our presentation from a blank slide template
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
sheetDict = xls_parsing_functions.read_xls_sheets(xlsFilename)
sheetsToExclude = set()
		#union the sheets we asked it to include with those that dont have proper columns
sheetsToExclude = sheetsToExclude.union(set(xls_parsing_functions.exclude_sheets_based_on_missing_columns(sheetDict)))
#remove all sheets to exlude from our dict
for sheet in sheetsToExclude:
	print sheetDict[sheet].columns
	print '-------'
	del sheetDict[sheet]

#Main loop: create a slide for every variant specified for inclusion by checking out all the sheets and rows of the excel sheet
for sheetName, df in sheetDict.items():
	#iterate over the rows of the sheet and if written as what we should include we include
	for index, row in df.iterrows():
		#Alert add an if statement for the 'export tab' the gcs will create 
		if True:
			slide = prs.slides.add_slide(blank_slide_layout)
			populate_slide(slide, row, udnId)
		#ALERT testing here is break statement 
		break
	#ALERT testing here is break statement
	break

saveFilename = udnId + '_curationSlides.pptx'
print 'finished, saving file as ', saveFilename
#filename is UDNID_slides.pptx
prs.save(saveFilename)
